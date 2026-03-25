from pptx import Presentation
from .base import BaseParser


class PPTParser(BaseParser):

    def parse(self, file_path: str) -> str:
        prs = Presentation(file_path)
        texts = []

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)

        return "\n".join(texts)
